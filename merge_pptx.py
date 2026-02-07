#!/usr/bin/env python3

# PPTX Merger

# Merge multiple PowerPoint (.pptx) files into a single presentation.
# Combines presentations in alphabetical order or custom order specified by the user.

# Author: Avijit Roy
# GitHub: https://github.com/heyavijitroy
# Repository: https://github.com/heyavijitroy/pptx-to-pdf-converter
# License: MIT

# Requirements:
# - python-pptx library
#   Install: pip install python-pptx

# Usage:
#     python merge_pptx.py
    
#     The script will merge all .pptx files in the current directory.
#     You can also specify files to merge in a specific order.


import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx library not installed")
    print("\nPlease install it using:")
    print("pip install python-pptx")
    sys.exit(1)

__version__ = "1.0.0"
__author__ = "Avijit Roy"


def merge_presentations(pptx_files, output_file, add_separator=False):
    """
    Merge multiple PowerPoint presentations into one.
    
    Args:
        pptx_files: List of Path objects for input .pptx files
        output_file: Path object for the output merged .pptx file
        add_separator: If True, adds a separator slide between presentations
    
    Returns:
        bool: True if merge was successful, False otherwise
    """
    try:
        # Create a new presentation (use first file as base)
        merged_prs = Presentation(pptx_files[0])
        total_slides = len(merged_prs.slides)
        print(f"  Base: {pptx_files[0].name} ({total_slides} slides)")
        
        # Add slides from remaining presentations
        for pptx_file in pptx_files[1:]:
            try:
                prs = Presentation(pptx_file)
                slide_count = len(prs.slides)
                
                # Optionally add separator slide
                if add_separator:
                    separator_slide = merged_prs.slides.add_slide(merged_prs.slide_layouts[6])  # Blank layout
                    # Add text to separator
                    textbox = separator_slide.shapes.add_textbox(
                        left=merged_prs.slide_width // 4,
                        top=merged_prs.slide_height // 2 - 500000,
                        width=merged_prs.slide_width // 2,
                        height=1000000
                    )
                    text_frame = textbox.text_frame
                    text_frame.text = f"📎 {pptx_file.stem}"
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = 440000  # 44pt
                        paragraph.font.bold = True
                    total_slides += 1
                
                # Copy all slides from current presentation
                for slide in prs.slides:
                    # Get the slide layout from the source
                    slide_layout = merged_prs.slide_layouts[0]  # Use title slide layout as default
                    
                    # Add new slide
                    new_slide = merged_prs.slides.add_slide(slide_layout)
                    
                    # Copy all shapes from the source slide
                    for shape in slide.shapes:
                        # Get the shape element
                        el = shape.element
                        # Create a copy and add to new slide
                        new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
                    
                    total_slides += 1
                
                print(f"  Added: {pptx_file.name} ({slide_count} slides)")
                
            except Exception as e:
                print(f"  ✗ Error adding {pptx_file.name}: {str(e)}")
                continue
        
        # Save the merged presentation
        merged_prs.save(str(output_file))
        print(f"\n✓ Successfully merged into: {output_file.name}")
        print(f"  Total slides: {total_slides}")
        return True
        
    except Exception as e:
        print(f"\n✗ Error during merge: {str(e)}")
        return False


def get_user_choice(pptx_files):
    """
    Ask user to choose merge order.
    
    Args:
        pptx_files: List of Path objects for available .pptx files
    
    Returns:
        tuple: (ordered_files, add_separator)
    """
    print("\nAvailable PowerPoint files:")
    for i, file in enumerate(pptx_files, 1):
        print(f"  {i}. {file.name}")
    
    print("\nOptions:")
    print("  1. Merge all files in alphabetical order")
    print("  2. Merge all files in custom order")
    print("  3. Select specific files to merge")
    
    while True:
        choice = input("\nEnter your choice (1-3): ").strip()
        
        if choice == "1":
            ordered_files = sorted(pptx_files)
            break
        elif choice == "2":
            print("\nEnter file numbers in desired order (space-separated):")
            print("Example: 3 1 2")
            order_input = input("> ").strip()
            try:
                indices = [int(x) - 1 for x in order_input.split()]
                ordered_files = [pptx_files[i] for i in indices]
                break
            except (ValueError, IndexError):
                print("Invalid input. Please try again.")
        elif choice == "3":
            print("\nEnter file numbers to merge (space-separated):")
            print("Example: 1 3 5")
            select_input = input("> ").strip()
            try:
                indices = [int(x) - 1 for x in select_input.split()]
                ordered_files = [pptx_files[i] for i in indices]
                break
            except (ValueError, IndexError):
                print("Invalid input. Please try again.")
        else:
            print("Invalid choice. Please enter 1, 2, or 3.")
    
    # Ask about separator slides
    while True:
        sep_choice = input("\nAdd separator slides between presentations? (y/n): ").strip().lower()
        if sep_choice in ['y', 'yes']:
            add_separator = True
            break
        elif sep_choice in ['n', 'no']:
            add_separator = False
            break
        else:
            print("Please enter 'y' or 'n'.")
    
    return ordered_files, add_separator


def main():
    """Main function to merge PowerPoint presentations."""
    
    # Print header
    print(f"\n{'='*60}")
    print(f"PPTX Merger v{__version__}")
    print(f"Author: {__author__}")
    print(f"{'='*60}\n")
    
    # Get the directory where the script is located
    script_dir = Path(__file__).parent.absolute()
    
    # Find all .pptx files in the directory
    pptx_files = list(script_dir.glob("*.pptx"))
    
    # Filter out temporary files (starting with ~$)
    pptx_files = [f for f in pptx_files if not f.name.startswith("~$")]
    
    if len(pptx_files) < 2:
        print("Error: Need at least 2 PowerPoint files to merge.")
        print("Please place multiple .pptx files in the same directory as this script.")
        return
    
    print(f"Found {len(pptx_files)} PowerPoint file(s) in the current directory.\n")
    
    # Get user's choice for merge order
    ordered_files, add_separator = get_user_choice(pptx_files)
    
    if len(ordered_files) < 2:
        print("\nError: Need to select at least 2 files to merge.")
        return
    
    # Get output filename
    default_output = "merged_presentation.pptx"
    output_name = input(f"\nEnter output filename (default: {default_output}): ").strip()
    if not output_name:
        output_name = default_output
    if not output_name.endswith('.pptx'):
        output_name += '.pptx'
    
    output_file = script_dir / output_name
    
    # Check if output file already exists
    if output_file.exists():
        overwrite = input(f"\n'{output_name}' already exists. Overwrite? (y/n): ").strip().lower()
        if overwrite not in ['y', 'yes']:
            print("Merge cancelled.")
            return
    
    # Perform the merge
    print(f"\n{'='*60}")
    print(f"Merging {len(ordered_files)} presentations...")
    print(f"{'='*60}\n")
    
    if merge_presentations(ordered_files, output_file, add_separator):
        print(f"\n{'='*60}")
        print("Merge completed successfully! 🎉")
        print(f"{'='*60}")
    else:
        print(f"\n{'='*60}")
        print("Merge failed. Please check the errors above.")
        print(f"{'='*60}")


if __name__ == "__main__":
    main()
